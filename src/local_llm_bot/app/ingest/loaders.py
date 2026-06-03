# Version: 1.1.3
# Changelog: 1.0.0 — AIStudio_733: first version header on loaders.py.
#            AIStudio_817: table-aware HTML/iXBRL extraction. Data <table> elements are
#            normalized to GFM markdown pipe-tables (colspan/rowspan expanded into a
#            rectangular grid, empty gutter rows/cols pruned, lone currency-symbol cells
#            merged into the adjacent value) and spliced into the text in document order
#            before get_text() flattening. Layout/decorative tables (the majority) fail the
#            data-table test and fall through to get_text() exactly as before. Shared,
#            format-agnostic back-half (_grid_to_markdown / _grid_is_data_table) is reusable
#            by future per-format grid extractors (xlsx/docx/pptx/pdf) — see
#            NOTES - AIStudio - Table Extraction Strategy - 2026-05-31.
#            1.1.0 — AIStudio_685 + AIStudio_885 (one cut, same get_text path, one re-ingest):
#            • 685 (table column-header binding): data <table> elements now serialize to
#              EXPLODED "row_label (column-key): value[unit]" statements instead of markdown
#              pipe-rows. 817 fixed row flattening but left stacked COLUMN headers detached
#              from values (Q5 misread JPM CET1 13.1% as the column-2 16.9%). Exploded form
#              binds every value to its full row+column coordinates so each fact is
#              self-contained — removes column ambiguity AND survives chunk splits. Applied
#              to ALL data tables (maximal self-containment lowers the model-capability bar
#              for small local models). New: _detect_header_band / _forward_fill_header_band /
#              _compose_column_keys / _grid_to_exploded / _table_to_text. The markdown
#              back-half (_grid_to_markdown / _table_to_markdown) is RETAINED for future
#              per-format extractors but no longer called on the HTML data-table path.
#            • 885 (iXBRL machine-layer strip): inline-XBRL filings carry a single
#              <ix:header> block (contexts/units/references + nested <ix:hidden> facts) whose
#              get_text() output is tag-soup (xbrli:context, explicitMember dimensions, bare
#              CIKs, dates) — observed corrupting whole chunks (e.g. CME chunk-14). Stripped
#              from the RAW markup before parsing via _strip_ixbrl_machine_layer (regex,
#              parser-independent; also a parse-speed win). Visible <ix:nonfraction>/
#              <ix:nonnumeric> body facts are untouched. NOTE: entity/year extraction in
#              pipeline.py parses the raw file independently, so this strip does not affect
#              the [Document: entity FY year] prefix (verify before re-ingest).
#            1.1.1 — AIStudio_685 validation fixes (real-grid stress test on JPM 2022):
#            • Row-label forward-fill: a sub-row whose col-0 is blank (e.g. "Diluted" under
#              "Net income:  Basic") now inherits the last non-empty col-0 label so the value
#              keeps its parent context. Runaway intra-cell whitespace collapsed.
#            • Column binding = EXACT-INDEX first, NEAREST-LEFT fallback: dense tables (the
#              dedicated capital table) bind by exact column index (unchanged — reads 13.1
#              Standardized / 13.8 Advanced correctly); only when the exact-index key is empty
#              (sparse/drifting summary tables whose year headers sit at columns the values
#              don't occupy) does a value fall back to the nearest non-empty header anchor
#              at-or-left. Cannot regress the dense case (it never reaches the fallback).
#            • Qualifier columns: a header anchor that is neither a date/period nor an entity
#              (e.g. "Change", "Capital ratio requirements") is folded into the key as a
#              qualifier ("(YoY Change)", "(regulatory minimum)") rather than emitted as a
#              bare period value — kills the Q5-class "1% revenue" misread while keeping the
#              datum (per Decision 1-A).
#            • Cosmetic: \xa0 → space; trailing footnote markers ((a)/(b)/(\u201cCET1\u201d)) trimmed
#              from composed keys/labels.
#            1.1.2 — AIStudio_685 gate polish: _grid_is_data_table now requires >=2 numeric cells
#            in the VALUE columns (cols>=1), not >=1 numeric anywhere. Rejects two-column
#            narrative/TOC/footnote tables (e.g. "OVERVIEW | 4", "* | See 'Compensation...'")
#            that previously passed on a stray page-number digit and produced anchorless,
#            value-less output. No effect on real data tables. Found via the 4-bank stress test
#            (Citi 26 anchorless tables — all narrative/TOC, confirmed no capital data lost).
#            1.1.3 — AIStudio_685 / AIStudio_880 European number-format fix: _cell_is_value now
#            recognises space/NBSP/narrow-NBSP/thin-space thousands separators ("27 254"),
#            apostrophe thousands ("1'234"), and comma-decimal ("10,29"); a date guard keeps
#            "31.12.2024"-style headers out of the value class. Before this, ESEF FR/DE/Nordic
#            numbers failed isdigit(), so their rows read as label-majority, got pulled into the
#            header band, and polluted anchors with leaked data (e.g. SocGen "2025, 27 254, ...").
#            sec_10k (comma-thousands) behaviour unchanged. SEB-class intra-word text shredding
#            is a SEPARATE still-open defect (AIStudio_880).
from __future__ import annotations

import fnmatch
import logging
import re
import warnings
from dataclasses import dataclass
from pathlib import Path

# from ..config import DEFAULT_XLSX_MAX_CELLS
from local_llm_bot.app.config import CONFIG

SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx", ".htm", ".html", ".xhtml"}

# Silence noisy PDF parsing logs/warnings (common with imperfect PDFs)
logging.getLogger("pypdf").setLevel(logging.ERROR)
warnings.filterwarnings("ignore", module="pypdf")


@dataclass(frozen=True)
class ExtractResult:
    ok: bool
    text: str = ""
    reason: str = ""  # empty if ok, otherwise a short failure label


def should_skip_filename(name: str) -> bool:
    return name.startswith("~$") or name == ".DS_Store"


def is_excluded(path: Path, patterns: list[str]) -> bool:
    s = str(path)
    return any(fnmatch.fnmatch(s, pat) for pat in patterns)


def extract_text(path: Path) -> ExtractResult:
    ext = path.suffix.lower()
    if ext in {".txt", ".md"}:
        return _extract_txt_md(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path, max_cells=CONFIG.ingest.xlsx_max_cells)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext in {".htm", ".html", ".xhtml"}:
        return _extract_html(path)
    return ExtractResult(ok=False, text="", reason="unsupported_ext")


def _extract_txt_md(path: Path) -> ExtractResult:
    try:
        t = path.read_text(encoding="utf-8", errors="ignore").strip()
        if not t:
            return ExtractResult(ok=False, text="", reason="empty")
        return ExtractResult(ok=True, text=t, reason="")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"read_error:{type(e).__name__}")


def _extract_docx(path: Path) -> ExtractResult:
    try:
        from docx import Document as DocxDocument  # type: ignore
    except Exception:
        return ExtractResult(ok=False, text="", reason="missing_dep:python-docx")

    try:
        doc = DocxDocument(str(path))
        parts: list[str] = []

        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        parts.append(t)

        text = "\n".join(parts).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        return ExtractResult(ok=True, text=text, reason="")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")


def _extract_pptx(path: Path) -> ExtractResult:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return ExtractResult(ok=False, text="", reason="missing_dep:python-pptx")

    try:
        prs = Presentation(str(path))
        parts: list[str] = []

        for slide in prs.slides:
            for shape in slide.shapes:
                # Text frames
                if hasattr(shape, "text") and shape.text:
                    t = str(shape.text).strip()
                    if t:
                        parts.append(t)

                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    tbl = shape.table
                    for row in tbl.rows:
                        for cell in row.cells:
                            t = (cell.text or "").strip()
                            if t:
                                parts.append(t)

        text = "\n".join(parts).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        return ExtractResult(ok=True, text=text, reason="")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")


def _extract_xlsx(path: Path, max_cells: int) -> ExtractResult:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return ExtractResult(ok=False, text="", reason="missing_dep:openpyxl")

    try:
        wb = openpyxl.load_workbook(filename=str(path), read_only=True, data_only=True)
        parts: list[str] = []
        scanned = 0
        partial = False

        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    scanned += 1
                    if scanned > max_cells:
                        partial = True
                        break
                    if v is None:
                        continue
                    s = str(v).strip()
                    if s:
                        parts.append(s)
                if partial:
                    break
            if partial:
                break

        wb.close()
        text = "\n".join(parts).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        if partial:
            return ExtractResult(ok=True, text=text, reason="partial:cell_cap")
        return ExtractResult(ok=True, text=text, reason="")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")


def _extract_pdf(path: Path) -> ExtractResult:
    """
    Extract PDF text with page boundary markers using pdfplumber.
    Each page is prefixed with [PAGE_N] so downstream chunking can
    store page numbers in Qdrant payload for click-to-source citation.
    Falls back to pypdf if pdfplumber is unavailable.
    """
    try:
        import pdfplumber  # type: ignore

        parts: list[str] = []
        with pdfplumber.open(str(path)) as pdf:
            if not pdf.pages:
                return ExtractResult(ok=False, text="", reason="empty")
            for page_num, page in enumerate(pdf.pages, 1):
                t = (page.extract_text() or "").strip()
                if t:
                    parts.append(f"[PAGE_{page_num}]\n{t}")

        text = "\n\n".join(parts).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        if text.lstrip().startswith("%PDF-"):
            return ExtractResult(ok=False, text="", reason="pdf_bytes_detected")
        return ExtractResult(ok=True, text=text, reason="")

    except ImportError:
        pass  # fall through to pypdf
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")

    # Fallback: pypdf (no page markers)
    try:
        from pypdf import PdfReader  # type: ignore
        from pypdf.errors import PdfReadError  # type: ignore
    except Exception:
        return ExtractResult(ok=False, text="", reason="missing_dep:pypdf")

    try:
        reader = PdfReader(str(path))
        if getattr(reader, "is_encrypted", False):
            return ExtractResult(ok=False, text="", reason="encrypted_pdf")

        parts = []
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(t)

        text = "\n".join(parts).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        if text.lstrip().startswith("%PDF-"):
            return ExtractResult(ok=False, text="", reason="pdf_bytes_detected")
        return ExtractResult(ok=True, text=text, reason="")
    except PdfReadError:
        return ExtractResult(ok=False, text="", reason="pdf_read_error")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")


# ── Table normalization (AIStudio_817) ───────────────────────────────────────
# Format-agnostic back-half: _grid_to_markdown / _grid_is_data_table operate on an
# abstract 2-D grid of cell strings and know nothing about HTML. The HTML-specific
# part is only _html_table_to_grid (colspan/rowspan expansion). Future per-format
# extractors (xlsx iter_rows, docx/pptx .rows[].cells[]) produce a grid and reuse the
# same back-half. See NOTES - AIStudio - Table Extraction Strategy - 2026-05-31 §3.

_CURRENCY_SYMBOLS = {"$", "€", "£", "¥"}


def _grid_put(matrix: list[list[str]], r: int, c: int, value: str) -> None:
    while len(matrix) <= r:
        matrix.append([])
    row = matrix[r]
    while len(row) <= c:
        row.append("")
    row[c] = value


def _grid_rectangularize(matrix: list[list[str]]) -> list[list[str]]:
    width = max((len(r) for r in matrix), default=0)
    for r in matrix:
        while len(r) < width:
            r.append("")
    return matrix


def _html_table_to_grid(table) -> list[list[str]]:
    """Expand an HTML <table> (with colspan/rowspan) into a rectangular cell matrix."""
    matrix: list[list[str]] = []
    carry: dict[tuple[int, int], str] = {}  # (row, col) -> text held by a rowspan
    for r, tr in enumerate(table.find_all("tr")):
        if len(matrix) <= r:
            matrix.append([])
        c = 0
        for cell in tr.find_all(["td", "th"], recursive=False):
            while (r, c) in carry:
                _grid_put(matrix, r, c, carry.pop((r, c)))
                c += 1
            style = (cell.get("style") or "").replace(" ", "").lower()
            text = "" if "display:none" in style else cell.get_text(" ", strip=True)
            colspan = int(cell.get("colspan", 1) or 1)
            rowspan = int(cell.get("rowspan", 1) or 1)
            for dc in range(colspan):
                val = text if dc == 0 else ""
                _grid_put(matrix, r, c + dc, val)
                for dr in range(1, rowspan):
                    carry[(r + dr, c + dc)] = val
            c += colspan
    return _grid_rectangularize(matrix)


def _grid_prune(grid: list[list[str]]) -> list[list[str]]:
    """Drop fully-empty rows and columns (spacer rows + gutter columns)."""
    grid = [r for r in grid if any(x.strip() for x in r)]
    if not grid:
        return grid
    keep = [c for c in range(len(grid[0])) if any(r[c].strip() for r in grid)]
    return [[r[c] for c in keep] for r in grid]


def _grid_merge_currency(grid: list[list[str]]) -> list[list[str]]:
    """Merge a lone currency-symbol cell into the next non-empty cell in its row."""
    out: list[list[str]] = []
    for r in grid:
        nr, i = [], 0
        while i < len(r):
            cur = r[i].strip()
            if cur in _CURRENCY_SYMBOLS and i + 1 < len(r) and r[i + 1].strip():
                nr.append(f"{cur}{r[i + 1].strip()}")
                i += 2
            else:
                nr.append(r[i])
                i += 1
        out.append(nr)
    return _grid_rectangularize(out)


def _grid_is_data_table(grid: list[list[str]]) -> bool:
    """True only for real data grids. Layout/decorative/narrative tables fail this and
    fall through to get_text().

    Gate: >=2x2 with >=4 non-empty cells AND >=2 numeric cells located in the VALUE
    columns (cols >=1). The value-column + count>=2 requirement (tightened 2026-06-03,
    AIStudio_685 v1.1.1) rejects two-column narrative/TOC/footnote tables — e.g.
    "OVERVIEW | 4" or "* | See 'Compensation Discussion'..." — which previously slipped
    through on a single stray page-number/marker digit and produced anchorless,
    value-less output. A genuine financial table always has numbers across its value
    columns, not one incidental digit."""
    if len(grid) < 2 or (grid and len(grid[0]) < 2):
        return False
    nonempty = sum(1 for r in grid for x in r if x.strip())
    numeric_vals = sum(1 for r in grid for x in r[1:] if _cell_is_value(x))
    return nonempty >= 4 and numeric_vals >= 2


def _grid_to_markdown(grid: list[list[str]]) -> str:
    """Render a grid as a GFM pipe-table (first row = header). Emits all rows."""
    head = grid[0]
    lines = ["| " + " | ".join(head) + " |",
             "| " + " | ".join(["---"] * len(head)) + " |"]
    for r in grid[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _table_to_markdown(table) -> str | None:
    """Normalize one HTML <table> to markdown, or None if it is not a data table.

    RETAINED (AIStudio_685) as the reusable markdown back-half for future per-format
    extractors; NO LONGER CALLED on the HTML data-table path — see _table_to_text.
    """
    grid = _grid_merge_currency(_grid_prune(_html_table_to_grid(table)))
    if not _grid_is_data_table(grid):
        return None
    return _grid_to_markdown(grid)


# ── iXBRL machine-layer stripping (AIStudio_885) ──────────────────────────────
# Inline-XBRL filings carry a single <ix:header> block (contexts, units, references,
# and a nested <ix:hidden> facts section) that is machine-only and never meant to
# render. get_text() would otherwise extract its contents as tag-soup:
# xbrli:context, xbrldi:explicitMember dimensions, bare CIKs, isolated dates
# (observed corrupting whole chunks, e.g. CME_Group chunk-14). We remove it from the
# RAW markup BEFORE BeautifulSoup parses — regex-based so it is parser-independent
# (immune to html.parser colon-namespace quirks) and shrinks the DOM before parse.
# Visible facts (<ix:nonfraction>/<ix:nonnumeric> in the body) are OUTSIDE <ix:header>
# and are left fully intact.
_IX_HEADER_RE = re.compile(r"<ix:header\b.*?</ix:header\s*>", re.I | re.S)
_IX_HIDDEN_RE = re.compile(r"<ix:hidden\b.*?</ix:hidden\s*>", re.I | re.S)


def _strip_ixbrl_machine_layer(html: str) -> str:
    html = _IX_HEADER_RE.sub("", html)
    html = _IX_HIDDEN_RE.sub("", html)  # belt-and-suspenders: if hidden sits outside header
    return html


# ── Table column-header binding / exploded serialization (AIStudio_685) ───────
# 817 serialized table ROWS correctly but left stacked COLUMN headers detached from
# their values, so a model reading "| CET1 capital ratio | 13.1 | 16.9 | ... |" had
# no in-row signal for which column is which and could pick the wrong one (Q5: read
# 13.1 vs 16.9). Fix: bind each value to its full row+column coordinates and emit one
# self-contained statement per cell — "row_label (column-key): value[unit]". Each
# fact then survives any chunk split AND carries zero column ambiguity. Applied to
# ALL data tables: maximal self-containment lowers the model-capability bar (a small
# local model cannot misread a fact that is already fully labeled).

_YEAR_RE = re.compile(r"^(19|20)\d{2}$")
# Day/month/year dates (31.12.2024, 31/12/2024, 12-31-2024) are HEADERS, not values —
# guard so they aren't misread as numbers once separators are stripped (AIStudio_685 v1.1.3).
_DATE_RE = re.compile(r"^\d{1,2}[./\-]\d{1,2}[./\-]\d{2,4}$")
# All separator/decoration chars stripped before the all-digits test. Includes European
# thousands separators: ASCII space, NBSP (\u00a0), NARROW NBSP (\u202f), THIN SPACE
# (\u2009) — ESEF filings (FR/DE/Nordic) group thousands with these (e.g. "27 254"), which
# previously failed isdigit(), mis-classified the row as a label, and polluted the header
# band/anchors. Also apostrophe thousands (Swiss "1'234"), comma/dot, currency, %, sign,
# parens.
_VALUE_STRIP = str.maketrans(
    "", "", " \u00a0\u202f\u2009\u2007,.%$€£¥'\u2019()+-"
)


def _cell_is_value(c: str) -> bool:
    """A numeric data value (amount / ratio / percent / European-formatted number), NOT a
    label, a bare year, or a date. Handles US (1,234.56), European (1.234,56 / 27 254 /
    27 254,5), Swiss (1'234) and parenthesised-negative ((9,256)) formats."""
    c = c.strip()
    if not c or _YEAR_RE.match(c) or _DATE_RE.match(c):
        return False
    if not any(ch.isdigit() for ch in c):  # a value must contain a digit ("%" alone is not)
        return False
    cleaned = c.translate(_VALUE_STRIP)
    return cleaned.isdigit()


def _cell_is_label(c: str) -> bool:
    """A column/row label: contains letters, or is a bare 4-digit year."""
    c = c.strip()
    return bool(re.search(r"[A-Za-z]", c)) or bool(_YEAR_RE.match(c))


_FOOTNOTE_RE = re.compile(r"\s*\((?:[a-z]{1,3}|\d{1,3}|[\u201c\u201d\"'][^)]*[\u201c\u201d\"'])\)\s*$")


def _clean_text(c: str) -> str:
    """Normalize a cell for keys/labels: \xa0 → space, collapse whitespace runs, strip."""
    return re.sub(r"\s+", " ", c.replace("\xa0", " ")).strip()


def _strip_footnotes(s: str) -> str:
    """Trim a trailing footnote marker — (a), (b), (12), (\u201cCET1\u201d) — from a key/label."""
    prev = None
    while prev != s:
        prev = s
        s = _FOOTNOTE_RE.sub("", s).strip()
    return s


def _detect_header_band(grid: list[list[str]]) -> int:
    """Return the index where the table BODY begins (= count of leading title/header rows).

    Rows above the first data row are the band. Three row kinds:
      • title/divider — <2 non-empty value cells (a bare "Year ended December 31..." or a
        section label like "Risk-based capital metrics:"); skipped INTO the band, never
        treated as a column header.
      • header — >=2 non-empty value cells, label-majority (text/year, not data values).
      • data — >=2 value cells, value-majority → STOP; this is the first body row.
    No <th>/<thead> survives gridding and SEC/ESEF filings use neither, so this is the
    heuristic path. (The earlier version broke on the leading title row and found zero
    headers — fixed: titles are skipped, not terminal.)
    """
    n = 0
    for row in grid:
        vals = [c for c in row[1:] if c.strip()]
        values = sum(1 for c in vals if _cell_is_value(c))
        labels = sum(1 for c in vals if _cell_is_label(c))
        if len(vals) >= 2 and values > labels:
            break  # first data row → body begins here
        n += 1
    return n


def _is_header_row(row: list[str]) -> bool:
    """A genuine column-header row has >=2 non-empty cells in the value columns (cols>=1).
    Col0-only rows (titles like "Year ended December 31...", dividers like "Risk-based
    capital metrics:") are NOT header rows — they must not be forward-filled across the
    row or composed into column keys."""
    return sum(1 for c in row[1:] if c.strip()) >= 2


def _forward_fill_header_band(grid: list[list[str]], n_header: int) -> list[list[str]]:
    """Return a COPY with empty cells in genuine header rows filled left-to-right.

    Colspan recovery: 817's _html_table_to_grid puts spanned header text only in the
    leftmost expanded cell; this carries it rightward (bounded by the next non-empty
    cell in the same row) so a stacked header like "Standardized" covers all its leaf
    columns. ONLY header rows are filled — filling a col0-only title/divider row would
    smear its label across every column and poison every key. Body untouched.
    """
    out = [list(r) for r in grid]
    for r in range(min(n_header, len(out))):
        if not _is_header_row(out[r]):
            continue  # title / divider — leave as-is
        last = ""
        for c in range(len(out[r])):
            if out[r][c].strip():
                last = out[r][c].strip()
            elif last:
                out[r][c] = last
    return out


def _anchor_keys(grid: list[list[str]], n_header: int) -> list[str]:
    """Ordered list of column-keys, one per LEAF HEADER ANCHOR, left-to-right.

    Anchor columns = the columns where the UNFILLED bottom header row (the leaf level)
    has content — i.e. the real distinct value columns. Each anchor's key is composed
    by reading the FORWARD-FILLED *header rows* top-to-bottom at that column, so stacked
    levels (Standardized / JPMorgan Chase & Co.) compose into one key while title/divider
    rows are skipped. This ordered list is what body values are zipped against (positional
    column index is unreliable in real filings — values drift off their headers — so we
    match by ORDER, not index).
    """
    if n_header <= 0:
        return []
    filled = _forward_fill_header_band(grid, n_header)
    header_rows = [r for r in range(n_header) if _is_header_row(grid[r])]
    if not header_rows:
        return []
    bottom = header_rows[-1]  # leaf header row
    keys: list[str] = []
    for c in range(1, len(grid[bottom])):
        if not grid[bottom][c].strip():
            continue  # not a leaf anchor column
        parts: list[str] = []
        for r in header_rows:
            if c < len(filled[r]):
                t = _strip_footnotes(_clean_text(filled[r][c]))
                if t and t not in parts:
                    parts.append(t)
        keys.append(", ".join(parts))
    return keys


def _row_values(row: list[str]) -> list[str]:
    """Ordered list of 'value[unit]' tokens in a body row (a '%' following a value across
    intervening empty gutter cells is attached)."""
    out: list[str] = []
    c = 1
    while c < len(row):
        cell = row[c].strip()
        if not _cell_is_value(cell):
            c += 1
            continue
        unit = ""
        j = c + 1
        while j < len(row) and not row[j].strip():
            j += 1
        if j < len(row) and row[j].strip() == "%":
            unit = "%"
        out.append(f"{cell}{unit}")
        c += 1
    return out


def _grid_to_exploded(grid: list[list[str]], n_header: int) -> str:
    """Emit one 'row_label (column-key): value' statement per body value, binding values
    to header anchors BY ORDER (the Nth value → the Nth anchor).

    - Row-label inheritance: a body row with a blank col-0 inherits the last non-empty
      col-0 label (sub-rows like "Diluted" keep their parent context).
    - Order-zip: robust to the positional drift that makes index/nearest-left binding
      mis-assign in real filings; the dense capital table (4 values → 4 leaf anchors) and
      the sparse summary table (3 values → 3 anchors) both bind correctly.
    - If a row has more values than anchors, the surplus emit with the row label only;
      if it has fewer, the first k zip and the rest are absent (correct — missing cells).
    """
    anchors = _anchor_keys(grid, n_header)
    lines: list[str] = []
    last_label = ""
    for row in grid[n_header:]:
        col0 = _strip_footnotes(_clean_text(row[0])) if row and row[0].strip() else ""
        if col0:
            last_label = col0
        row_label = col0 or last_label
        values = _row_values(row)
        for i, val in enumerate(values):
            key = anchors[i] if i < len(anchors) else ""
            if row_label and key:
                lines.append(f"{row_label} ({key}): {val}")
            elif key:
                lines.append(f"{key}: {val}")
            elif row_label:
                lines.append(f"{row_label}: {val}")
            else:
                lines.append(val)
    return "\n".join(lines)


def _table_to_text(table) -> str | None:
    """Normalize one HTML <table> to EXPLODED label:value statements (AIStudio_685),
    or None if it is not a data table. Detect header band → compose ordered leaf-anchor
    keys → emit one statement per body value, zipped to anchors by order.
    """
    grid = _grid_merge_currency(_grid_prune(_html_table_to_grid(table)))
    if not _grid_is_data_table(grid):
        return None
    n_header = _detect_header_band(grid)
    return _grid_to_exploded(grid, n_header)


def _extract_html(path: Path) -> ExtractResult:
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except Exception:
        return ExtractResult(ok=False, text="", reason="missing_dep:beautifulsoup4")

    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        # AIStudio_885: strip the iXBRL machine layer (<ix:header>/<ix:hidden>) from the
        # raw markup before parsing — removes the context/dimension/CIK tag-soup that
        # would otherwise flatten into chunks; parser-independent and shrinks the DOM.
        raw = _strip_ixbrl_machine_layer(raw)
        soup = BeautifulSoup(raw, "html.parser")
        # Remove scripts, styles, nav boilerplate
        for tag in soup(["script", "style", "nav", "header", "footer", "meta", "link"]):
            tag.decompose()
        # AIStudio_817 + AIStudio_685: convert data tables IN DOCUMENT ORDER before the
        # get_text() flatten. Each <table> that normalizes to a data grid is replaced
        # in-place by EXPLODED "row_label (column-key): value" statements (685; wrapped
        # in blank lines so chunking sees one atomic block); layout/decorative tables
        # return None and flatten as before. The parent-None guard skips tables already
        # detached by an outer replacement (nested-table case).
        for table in soup.find_all("table"):
            if table.parent is None:
                continue
            md = _table_to_text(table)
            if md:
                table.replace_with("\n\n" + md + "\n\n")
        text = soup.get_text(separator="\n").strip()
        # Collapse excessive blank lines
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        if not text:
            return ExtractResult(ok=False, text="", reason="empty")
        return ExtractResult(ok=True, text=text, reason="")
    except Exception as e:
        return ExtractResult(ok=False, text="", reason=f"parse_error:{type(e).__name__}")


# Backward-compatible helper if you already import load_document elsewhere:
@dataclass(frozen=True)
class Document:
    doc_id: str
    source_path: str
    text: str


def load_document(path: Path) -> Document | None:
    if not path.is_file():
        return None
    if should_skip_filename(path.name):
        return None

    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        return None

    res = extract_text(path)
    if not res.ok or not res.text.strip():
        return None

    return Document(doc_id=str(path.resolve()), source_path=str(path), text=res.text)
